import streamlit as st
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os

st.title("📄 PDF를 PPTX로 변환하기")
st.write("PDF 파일을 업로드하면 각 페이지를 이미지로 변환하여 PPTX 파일로 만들어드립니다!")

uploaded_file = st.file_uploader("PDF 파일을 업로드하세요", type=["pdf"])

if uploaded_file is not None:
    # 임시 폴더 생성
    with tempfile.TemporaryDirectory() as tmp_dir:
        # 업로드된 PDF 파일 저장
        pdf_data = uploaded_file.read()

        # PDF를 이미지로 변환
        st.write("📷 PDF를 이미지로 변환 중입니다...")
        try:
            poppler_path = os.getenv("POPPLER_PATH", "/usr/bin")  # 환경 변수 또는 기본 경로 사용
            images = convert_from_bytes(pdf_data, poppler_path=poppler_path)
        except Exception as e:
            st.error(f"PDF를 이미지로 변환하는 중 오류가 발생했습니다: {e}")
            st.stop()

        # PPTX 생성
        st.write("📽 PPTX 파일을 생성 중입니다...")
        presentation = Presentation()

        for i, image in enumerate(images):
            # 이미지 파일 임시 저장
            image_path = f"{tmp_dir}/page_{i + 1}.jpg"
            image.save(image_path, "JPEG")

            # 슬라이드 추가 및 이미지 삽입
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)

        # PPTX 파일 저장
        pptx_path = f"{tmp_dir}/converted_presentation.pptx"
        presentation.save(pptx_path)

        # 변환된 PPTX 다운로드 링크 제공
        with open(pptx_path, "rb") as f:
            st.download_button(
                label="📥 PPTX 파일 다운로드",
                data=f,
                file_name="converted_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        st.success("✅ 변환이 완료되었습니다! PPTX 파일을 다운로드하세요.")
